"""Generate latex/final_presentation/timeline.pptx.

Layout: positions are taken directly from latex/final_presentation/Timeline_Component_Deck.pptx
(a 1:1 px->EMU (9525 EMU/px) export of the original 1920x1080 'Timeline Component.dc.html'
canvas), so the timeline track sits at the same relative height (~598/1080 = middle of the
slide) as the live component, with milestone content alternating above/below it.

One deliberate simplification vs. the live component: top-anchored milestones there are
bottom-aligned (CSS flex-end, hugging the track) with height depending on bullet count;
here every top-anchored block is top-aligned at the same fixed y (mirroring how
bottom-anchored blocks already behave in both the component and the baseline export), since
pptx has no flexbox equivalent and content is now variable-length bulleted lists rather than
the baseline's fixed single-line body.

Transition: PowerPoint has no way to scope a transition to part of a slide; Morph is a
whole-slide effect that auto-matches similar-looking shapes across slides, which is why an
earlier attempt at Morph "morphed everything" (headline/body text got matched and slid
around). The fix is PowerPoint's documented '!!Name' convention: shapes sharing an identical
'!!'-prefixed name across two slides are force-matched by Morph (and interpolate position/
size/color); shapes without a shared '!!' name are left unmatched and simply cross-fade
instead. Only the timeline track, progress fill, and the 9 dots get stable '!!' names here;
every other shape (pill/headline/bullets/image placeholder) gets a unique per-slide name so
Morph can never accidentally match it to its counterpart on the next slide.
Sources (see chat): MS Support "Morph transition: tips and tricks", Indezine/Office-Watch on
'!!' exclamation-named objects, and the MS-PPTX open spec on CT_MorphTransition (namespace
http://schemas.microsoft.com/office/powerpoint/2015/09/main, option="byObject").
"""
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# Palette/font below are lifted from DPaDS_final.pptx's theme ("SVSTemplate2023",
# the Uni Hamburg template) rather than eyeballed, so this deck reads as the same
# deliverable family as the final presentation. Geometry/positions elsewhere in this
# file are untouched -- only these constants, add_divider_line(), and add_logo() (both
# new, additive shapes) changed; see README.md for what "design-only" means here.
ACCENT = RGBColor(0x32, 0x57, 0x86)        # DPaDS_final theme dk2/tx2 -- was 0x245C9E
BULLET_RED = RGBColor(0xC1, 0x21, 0x2A)    # DPaDS_final theme accent2 -- bullet markers there are red
PILL_BG = RGBColor(0xC9, 0xE1, 0xFF)       # DPaDS_final theme accent5 -- was 0xEAF1FB
BODY_GREY = RGBColor(0x4B, 0x55, 0x63)
MUTED_GREY = RGBColor(0x9C, 0xA3, 0xAF)
INK = RGBColor(0x1F, 0x20, 0x23)
TITLE_GREY = RGBColor(0x80, 0x80, 0x80)    # DPaDS_final titleStyle: tx1 (black), lumMod/lumOff 50%/50%
SUBTITLE_GREY = RGBColor(0x6B, 0x72, 0x80)
TRACK_GREY = RGBColor(0xE2, 0xE2, 0xE2)
DOT_INACTIVE = RGBColor(0xD1, 0xD5, 0xDB)
PLACEHOLDER_BG = RGBColor(0xED, 0xEE, 0xF0)
PLACEHOLDER_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
DIVIDER_GREY = RGBColor(0xB3, 0xB3, 0xB3)  # DPaDS_final's title-underline divider color
FONT = "Calibri"                           # DPaDS_final theme major/minor font -- was "Arial"

LOGO_PATH = Path(__file__).parent / "assets" / "uhh_logo.png"  # extracted from DPaDS_final.pptx (image1.png)
LOGO_ASPECT = 268 / 111  # native px size of the extracted logo

MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
P159_NS = "http://schemas.microsoft.com/office/powerpoint/2015/09/main"

MILESTONES = [
    {"type": "text", "date": "09.04.2026", "headline": "Project Kick-off", "bullets": [
        "Master Project Seminar Data Protection and Data Security begins"
    ]},
    {"type": "text", "date": "16.04.2026", "headline": "Schedule & Topic Distribution", "bullets": [
        "Presentation of project schedule and topic distribution"
    ]},
    {"type": "text", "date": "30.04.2026", "headline": "Small Presentations on Potential Topics", "bullets": [
        "DFL.tex - Security and Privacy Concerns in Federated Learning: An Introduction",
        "SMPC.tex - Introduction to Secure Multi-Party Computation"
    ]},
    {"type": "text", "date": "~21.05.2026", "headline": "Lock-in: FL Trilemma", "bullets": [
        "Theoretical understanding of the problem space",
        "Robust aggregation landscape reviewed: FLTrust, Krum, Trimmed Mean"
    ]},
    {"type": "image", "date": "18.06.2026", "headline": "Interim Presentation 1", "imageLabel": "INTERIM PRESENTATION 1 - SLIDES", "bullets": [
        "First code foundation building up (MNIST, Flower)"
    ]},
    {"type": "image", "date": "25.06.2026", "headline": "Interim Presentation 2", "imageLabel": "INTERIM PRESENTATION 2 - SLIDES", "bullets": [
        "Mechanisms implemented albeit partly flawed (DP, FLTrust, TopK)",
        "Testing extremely slow due to missing resources"
    ]},
    {"type": "image", "date": "09.07.2026", "headline": "Unlock of Computational Resources", "imageLabel": "ACCURACY / RUNTIME CHARTS", "bullets": [
        "Larger scale testing unlocked",
        "GPU access and incrementally stable code foundation allows for CIFAR-10 testing",
        "First results for large runs are coming through"
    ]},
    {"type": "text", "date": "16.07.2026", "headline": "Final Report Submission", "bullets": [
        "Written report on the Privacy-Robustness-Performance Trilemma due"
    ]},
    {"type": "image", "date": "20.07.2026", "headline": "Final Presentation", "imageLabel": "FINAL PRESENTATION - SLIDES/RESULTS", "bullets": [
        "Closing presentation of findings"
    ]},
]

# --- layout constants, copied 1:1 from Timeline_Component_Deck.pptx (EMU) ---
CANVAS_W, CANVAS_H = Emu(18288000), Emu(10287000)
TITLE = dict(x=1143000, y=666750, w=17602200, h=645319)
SUBTITLE = dict(x=1143000, y=1388194, w=17602200, h=347663)

PILL_H = 376163
PILL_W = 1650000
HEADLINE_OFFSET = 547502
HEADLINE_H = 478631
TEXT_BODY_OFFSET = 554757
TEXT_BODY_W = 11525250
IMG_CAPTION_OFFSET = 535782
# Widened past the baseline's 5451574 EMU (572px, sized for short single-line headlines like
# "Product Launch"): "Unlock of Computational Resources" alone measures ~689px at 30pt bold
# (measured with Liberation Sans, metric-compatible with Arial) and was wrapping to a second
# line that collided with the bullets below it. 1050px clears every headline/bullet in this
# deck with margin, while still leaving the box comfortably inside the 1920px canvas.
IMG_CAPTION_W = 1050 * 9525

TOP_ANCHOR_Y = 2190750      # 230px - component's top-container top; simplified to top-align
BOTTOM_ANCHOR_Y = 6667500   # 700px - component's bottom-container top (matches baseline exactly)

TEXT_X = 1524000        # 160px
IMAGE_TEXT_X = 5143500   # 540px (160 + 340 box + 40 gap)
IMAGE_BOX_X = 1524000    # 160px
IMAGE_BOX_W = 3238500    # 340px
IMAGE_BOX_H = 2095500    # 220px

TRACK_X, TRACK_Y, TRACK_W, TRACK_H = 1905000, 5695838, 14478000, 37951  # 200px,598px,1520px,4px
DOT_ANCHOR_Y_PX = 600
LABEL_Y = 6019726    # 632px
LEFT_MARGIN_PX, RIGHT_MARGIN_PX = 200, 1720

# New, additive design elements (see top-of-file comment) -- not part of the original
# Timeline_Component_Deck.pptx layout baseline above.
DIVIDER_X_LEFT = 1143000               # matches TITLE/SUBTITLE left edge
DIVIDER_X_RIGHT = CANVAS_W - 1143000   # symmetric right margin
DIVIDER_Y = 1905000                    # 200px -- modest gap below the subtitle block

LOGO_W = 2857500   # 300px
LOGO_H = round(LOGO_W / LOGO_ASPECT)
LOGO_X = CANVAS_W - LOGO_W - 952500   # 100px right margin -- top-right, clear of the
                                        # left-aligned title/subtitle/author text
LOGO_Y = 571500     # 60px top margin

prs = Presentation()
prs.slide_width = CANVAS_W
prs.slide_height = CANVAS_H
BLANK = prs.slide_layouts[6]
N = len(MILESTONES)


def add_fade_transition(slide):
    sld = slide._element
    trans = etree.SubElement(sld, f"{{{P_NS}}}transition")
    trans.set("spd", "slow")
    etree.SubElement(trans, f"{{{P_NS}}}fade")
    anchor = sld.find(f"{{{P_NS}}}clrMapOvr")
    if anchor is None:
        anchor = sld.find(f"{{{P_NS}}}cSld")
    anchor.addnext(trans)


def add_morph_transition(slide):
    sld = slide._element
    alt = etree.SubElement(sld, f"{{{MC_NS}}}AlternateContent", nsmap={"mc": MC_NS})
    choice = etree.SubElement(alt, f"{{{MC_NS}}}Choice", nsmap={"p159": P159_NS})
    choice.set("Requires", "p159")
    trans = etree.SubElement(choice, f"{{{P_NS}}}transition")
    trans.set("spd", "slow")
    morph = etree.SubElement(trans, f"{{{P159_NS}}}morph")
    morph.set("option", "byObject")

    fallback = etree.SubElement(alt, f"{{{MC_NS}}}Fallback")
    ftrans = etree.SubElement(fallback, f"{{{P_NS}}}transition")
    ftrans.set("spd", "slow")
    etree.SubElement(ftrans, f"{{{P_NS}}}fade")

    anchor = sld.find(f"{{{P_NS}}}clrMapOvr")
    if anchor is None:
        anchor = sld.find(f"{{{P_NS}}}cSld")
    anchor.addnext(alt)


def add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT, name=None, font=FONT):
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return box


def add_pill(slide, left, top, text, name):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(PILL_W), Emu(PILL_H))
    shp.name = name
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = PILL_BG
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = ACCENT
    return shp


def add_bullets(slide, left, top, width, bullets, size, name):
    height = 376163 * max(1, len(bullets))  # generous, never clips
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = "-  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.name = FONT
        r1.font.color.rgb = BULLET_RED   # DPaDS_final bodyStyle lvl1: red (accent2) bullet marker
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(size)
        r2.font.name = FONT
        r2.font.color.rgb = ACCENT       # DPaDS_final bodyStyle lvl1: blue (tx2) body text
    return box


def add_image_placeholder(slide, left, top, label, name):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(IMAGE_BOX_W), Emu(IMAGE_BOX_H))
    shp.name = name
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = PLACEHOLDER_BG
    shp.line.color.rgb = PLACEHOLDER_BORDER
    shp.line.width = Pt(1.5)
    shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(120000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(11)
    r.font.name = "Consolas"
    r.font.color.rgb = MUTED_GREY
    return shp


def add_divider_line(slide, name):
    """
    Thin grey line under the deck title/subtitle, mirroring DPaDS_final's
    title-underline divider (slideMaster1.xml's "Gerade Verbindung 14").

    Identical position/size/color on every slide, so it's harmless even if Morph
    auto-matches it across the milestone slides (interpolating between two
    identical states is a no-op) -- deliberately NOT given a `!!` name, since it
    isn't one of the timeline's own animated elements (see README.md).
    """
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(DIVIDER_X_LEFT), Emu(DIVIDER_Y),
                                       Emu(DIVIDER_X_RIGHT), Emu(DIVIDER_Y))
    line.name = name
    line.line.color.rgb = DIVIDER_GREY
    line.line.width = Pt(1.5)
    line.shadow.inherit = False
    return line


def add_logo(slide, name):
    """
    Uni Hamburg logo, matching DPaDS_final's title-slide branding (extracted from
    its slideLayout1.xml as assets/uhh_logo.png). DPaDS_final places it top-left;
    placed top-right here instead so it doesn't collide with this deck's own
    title/subtitle/author text, which all start at the left margin. Title slide
    only, matching DPaDS_final's own pattern (its content-slide layout has no logo).
    """
    pic = slide.shapes.add_picture(str(LOGO_PATH), Emu(LOGO_X), Emu(LOGO_Y), Emu(LOGO_W), Emu(LOGO_H))
    pic.name = name
    return pic


def dot_geometry(j, active_idx):
    spacing = (RIGHT_MARGIN_PX - LEFT_MARGIN_PX) / (N - 1)
    cx = LEFT_MARGIN_PX + j * spacing
    active = j == active_idx
    size_px = 28 if active else 16
    left_px = cx - size_px / 2
    top_px = DOT_ANCHOR_Y_PX - size_px / 2
    return cx, left_px, top_px, size_px, active


def add_timeline(slide, active_idx):
    track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(TRACK_X), Emu(TRACK_Y), Emu(TRACK_W), Emu(TRACK_H))
    track.name = "!!tl_track"
    track.adjustments[0] = 0.5
    track.fill.solid()
    track.fill.fore_color.rgb = TRACK_GREY
    track.line.fill.background()
    track.shadow.inherit = False

    fill_w = max(1, round(active_idx / (N - 1) * TRACK_W))
    fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(TRACK_X), Emu(TRACK_Y), Emu(fill_w), Emu(TRACK_H))
    fill.name = "!!tl_fill"
    fill.adjustments[0] = 0.5
    fill.fill.solid()
    fill.fill.fore_color.rgb = ACCENT
    fill.line.fill.background()
    fill.shadow.inherit = False

    for j, m in enumerate(MILESTONES):
        cx, left_px, top_px, size_px, active = dot_geometry(j, active_idx)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(round(left_px * 9525)), Emu(round(top_px * 9525)),
                                      Emu(round(size_px * 9525)), Emu(round(size_px * 9525)))
        dot.name = f"!!tl_dot_{j + 1}"
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT if active else DOT_INACTIVE
        dot.line.fill.background()
        dot.shadow.inherit = False

        add_text(slide, round((cx - 60) * 9525), LABEL_Y, round(120 * 9525), 300038,
                  m["date"], 18, INK if active else MUTED_GREY, bold=active, align=PP_ALIGN.CENTER,
                  name=f"lbl_{j + 1}_s{active_idx}")


# --- Title slide ---
title_slide = prs.slides.add_slide(BLANK)
add_text(title_slide, TITLE["x"], TITLE["y"], TITLE["w"], TITLE["h"], "FL-Trilemma: Project Timeline", 42, TITLE_GREY, bold=True)
add_text(title_slide, SUBTITLE["x"], SUBTITLE["y"], SUBTITLE["w"], SUBTITLE["h"], "From kick-off to final presentation", 21, SUBTITLE_GREY)
add_text(title_slide, TEXT_X, 3000000, 16764000, 400000, "Jonas Müller  ·  Rebekka Schnoor", 18, BODY_GREY)
add_divider_line(title_slide, name="divider_title")
add_logo(title_slide, name="logo_title")
add_fade_transition(title_slide)

# --- Milestone slides ---
for idx, m in enumerate(MILESTONES):
    slide = prs.slides.add_slide(BLANK)
    add_text(slide, TITLE["x"], TITLE["y"], TITLE["w"], TITLE["h"], "FL-Trilemma: Project Timeline", 42, TITLE_GREY, bold=True,
              name=f"title_m{idx}")
    add_text(slide, SUBTITLE["x"], SUBTITLE["y"], SUBTITLE["w"], SUBTITLE["h"], "From kick-off to final presentation", 21, SUBTITLE_GREY,
              name=f"subtitle_m{idx}")
    add_divider_line(slide, name=f"divider_m{idx}")

    anchor_y = TOP_ANCHOR_Y if idx % 2 == 0 else BOTTOM_ANCHOR_Y

    if m["type"] == "image":
        text_x = IMAGE_TEXT_X
        add_image_placeholder(slide, IMAGE_BOX_X, anchor_y, m["imageLabel"], name=f"content_placeholder_m{idx}")
        caption_offset, caption_w, caption_sz = IMG_CAPTION_OFFSET, IMG_CAPTION_W, 19.5
    else:
        text_x = TEXT_X
        caption_offset, caption_w, caption_sz = TEXT_BODY_OFFSET, TEXT_BODY_W, 21

    add_pill(slide, text_x, anchor_y, m["date"], name=f"content_pill_m{idx}")
    add_text(slide, text_x, anchor_y + HEADLINE_OFFSET, caption_w, HEADLINE_H, m["headline"], 30, INK, bold=True,
              name=f"content_headline_m{idx}")
    add_bullets(slide, text_x, anchor_y + HEADLINE_OFFSET + caption_offset, caption_w, m["bullets"], caption_sz,
                name=f"content_bullets_m{idx}")

    add_timeline(slide, idx)
    add_morph_transition(slide)

out_path = Path(__file__).parent / "timeline_base.pptx"
prs.save(out_path)
print("saved", out_path)
